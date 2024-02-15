# 自用
# 导入pandas模块 注意安装pandas和openyxl，（注意pandas已更名为pyarrow）
import numpy as np
import time
import pandas as pd

#注意安装pillow库
import win32api
import win32con
import win32gui
import os.path as op
import time
from PIL import Image
import sys
day=[0,0,0,0,0,0,0]
'''
#pip install fonttools
from fontTools.ttLib import TTFont

# 字体文件路径
font_path = "table/font/Aldrich-Regular.ttf"

# 加载字体文件
font = TTFont(font_path)

#安装字体
font.save("C:/Windows/Fonts/Aldrich-Regular.ttf")
'''
# 定义day, ts, te, sub为长度为7的列表，每个元素为一个空数组
sub = [[0] * 100] * 100
te = [[0] * 100] * 100
ts = [[0] * 100] * 100
# 定义一个函数，用于获取Excel文件的列数
def get_col_num(ws):
    # 返回ws的列数
    return ws.shape[1]

# 调用time模块的ctime()函数，返回当前的日期和时间，格式为字符串
ti = time.ctime()

#判断今天是星期几(0-6)
from datetime import datetime, date
dayOfWeek = datetime.today().weekday()

#导入python-pptx模块
from pptx import Presentation

# 循环遍历7个工作表
for i in range(7):
    # 读取Excel文件，指定sheet_name为i
    day[i] = pd.read_excel("ClassSchedule.xlsx", sheet_name=i)

    # 获取行数，传入day[i]作为参数
    
    rows,columns = day[i].shape

    # 循环遍历每一行
    for j in range(rows):
        
        # 读取第j+1行第2列的单元格的内容，赋值给ts[i][j]
        ts[i][j] = day[i].iloc[j, 1]
        
        # 读取第j+1行第3列的单元格的内容，赋值给te[i][j]
        te[i][j] = day[i].iloc[j, 2]
        
        # 读取第j+1行第4列的单元格的内容，赋值给sub[i][j]
        sub[i][j] = day[i].iloc[j, 3]

    # 打开现有的PPTX文件
    ppt = Presentation("PPT.pptx")
    
    from pptx import Presentation


    # 打开PPT文件
    ppt = Presentation("PPT.pptx")

    # 获取第一张幻灯片
    slide = ppt.slides[0]
    # 设置文本框的位置和尺寸
    left, top, width, height = 5549900, 1819982, 5425208, 2853904

    # 遍历幻灯片中的所有形状
    for shape in slide.shapes:
        # 判断形状是否为文本框
        if shape.has_text_frame:
            # 获取文本框的位置和尺寸
            text_box_left, text_box_top, text_box_width, text_box_height = shape.left, shape.top, shape.width, shape.height

            # 判断文本框的位置和尺寸是否与目标一致
            if (left, top, width, height) == (text_box_left, text_box_top, text_box_width, text_box_height):
                # 修改文本框内容为图片
                shape.text = ""
                picture = slide.shapes.add_picture("996.jpg", left, top, width, height)
                break

                # 将图片下移80个图层
                for i in range(80):
                    picture.z_order(picture.z_order - 1)
                
                break

    # 设置文本框的位置和尺寸
    left, top, width, height = 2809874, 3106057, 3286125, 1377329

    # 遍历幻灯片中的所有形状
    for shape in slide.shapes:
        # 判断形状是否为文本框
        if shape.has_text_frame:
            # 获取文本框的位置和尺寸
            text_box_left, text_box_top, text_box_width, text_box_height = shape.left, shape.top, shape.width, shape.height

            # 判断文本框的位置和尺寸是否与目标一致
            if (left, top, width, height) == (text_box_left, text_box_top, text_box_width, text_box_height):
                # 修改文本框内容
                shape.text = "xixi"
                break

    # 设置文本框的位置和尺寸
    left, top, width, height = 3288506, 2359820, 3064210, 746238

    # 遍历幻灯片中的所有形状
    for shape in slide.shapes:
        # 判断形状是否为文本框
        if shape.has_text_frame:
            # 获取文本框的位置和尺寸
            text_box_left, text_box_top, text_box_width, text_box_height = shape.left, shape.top, shape.width, shape.height

            # 判断文本框的位置和尺寸是否与目标一致
            if (left, top, width, height) == (text_box_left, text_box_top, text_box_width, text_box_height):
                # 修改文本框内容
                shape.text = "上课"
                break

    # 设置文本框的位置和尺寸
    left, top, width, height =11867892, 185766, 1613877,1054386

    # 遍历幻灯片中的所有形状
    for shape in slide.shapes:
        # 判断形状是否为文本框
        if shape.has_text_frame:
            # 获取文本框的位置和尺寸
            text_box_left, text_box_top, text_box_width, text_box_height = shape.left, shape.top, shape.width, shape.height

            # 判断文本框的位置和尺寸是否与目标一致
            if (left, top, width, height) == (text_box_left, text_box_top, text_box_width, text_box_height):
                # 修改文本框内容
                shape.text = "AA"
                break

    # 设置文本框的位置和尺寸
    left, top, width, height = 11492576, 1261632, 1613877, 1054386

    # 遍历幻灯片中的所有形状
    for shape in slide.shapes:
        # 判断形状是否为文本框
        if shape.has_text_frame:
            # 获取文本框的位置和尺寸
            text_box_left, text_box_top, text_box_width, text_box_height = shape.left, shape.top, shape.width, shape.height

            # 判断文本框的位置和尺寸是否与目标一致
            if (left, top, width, height) == (text_box_left, text_box_top, text_box_width, text_box_height):
                # 修改文本框内容
                shape.text = "222"
                break

    # 设置文本框的位置和尺寸
    left, top, width, height = 11118590, 2340967, 1613877, 1054386

    # 遍历幻灯片中的所有形状
    for shape in slide.shapes:
        # 判断形状是否为文本框
        if shape.has_text_frame:
            # 获取文本框的位置和尺寸
            text_box_left, text_box_top, text_box_width, text_box_height = shape.left, shape.top, shape.width, shape.height

            # 判断文本框的位置和尺寸是否与目标一致
            if (left, top, width, height) == (text_box_left, text_box_top, text_box_width, text_box_height):
                # 修改文本框内容
                shape.text = "333"
                break

    # 设置文本框的位置和尺寸
    left, top, width, height =10732921,3421380, 1613877, 1054386

    # 遍历幻灯片中的所有形状
    for shape in slide.shapes:
        # 判断形状是否为文本框
        if shape.has_text_frame:
            # 获取文本框的位置和尺寸
            text_box_left, text_box_top, text_box_width, text_box_height = shape.left, shape.top, shape.width, shape.height
    
            # 判断文本框的位置和尺寸是否与目标一致
            if (left, top, width, height) == (text_box_left, text_box_top, text_box_width, text_box_height):
                # 修改文本框内容
                shape.text = "444"
                break
    # 设置文本框的位置和尺寸
    left, top, width, height =10349752,4494173, 1613877, 1054386

    # 遍历幻灯片中的所有形状
    for shape in slide.shapes:
        # 判断形状是否为文本框
        if shape.has_text_frame:
            # 获取文本框的位置和尺寸
            text_box_left, text_box_top, text_box_width, text_box_height = shape.left, shape.top, shape.width, shape.height

            # 判断文本框的位置和尺寸是否与目标一致
            if (left, top, width, height) == (text_box_left, text_box_top, text_box_width, text_box_height):
                # 修改文本框内容
                shape.text = "555"
                break

    # 设置文本框的位置和尺寸
    left, top, width, height =9977901,5581128, 1613877, 1054386

    # 遍历幻灯片中的所有形状
    for shape in slide.shapes:
        # 判断形状是否为文本框
        if shape.has_text_frame:
            # 获取文本框的位置和尺寸
            text_box_left, text_box_top, text_box_width, text_box_height = shape.left, shape.top, shape.width, shape.height

            # 判断文本框的位置和尺寸是否与目标一致
            if (left, top, width, height) == (text_box_left, text_box_top, text_box_width, text_box_height):
                # 修改文本框内容
                shape.text = "666"
                break
    # 设置文本框的位置和尺寸
    left, top, width, height =9602816,6659737, 1613877, 1054386

    # 遍历幻灯片中的所有形状
    for shape in slide.shapes:
        # 判断形状是否为文本框
        if shape.has_text_frame:
            # 获取文本框的位置和尺寸
            text_box_left, text_box_top, text_box_width, text_box_height = shape.left, shape.top, shape.width, shape.height

            # 判断文本框的位置和尺寸是否与目标一致
            if (left, top, width, height) == (text_box_left, text_box_top, text_box_width, text_box_height):
                # 修改文本框内容
                shape.text = "777"
                break  


    # 保存修改后的PPT文件
    ppt.save("PPT2.pptx")
    # 打开PPT文件
    prs = Presentation("PPT2.pptx")
    import os
    # 遍历幻灯片中的所有形状
    for i, slide in enumerate(prs.slides):
        for j, shape in enumerate(slide.shapes):
            # 判断形状是否为文本框或图片
            if shape.has_text_frame or shape.has_image:
                # 获取形状的位置和尺寸
                left, top, width, height = shape.left, shape.top, shape.width, shape.height

                # 设置截图区域
                crop_area = (left, top, left + width, top + height)

                # 保存截图为图片
                image_file = f"images/slide{i+1}_shape{j+1}.png"
                slide.export(image_file, "PNG", crop_area)

        
#导入datetime 模块
import datetime
#无限循环
from itertools import cycle
for i in cycle((None,)):
    
    # 导入datetime模块
    import datetime

    # 获取当前的日期和时间
    now = datetime.datetime.now()

    # 打印当前的日期和时间，使用默认的格式
    print(now.strftime("%H:%M:%S"))
    
    # 调用time模块的sleep()函数，暂停程序执行0.75秒
    time.sleep(0.75)
    
    
