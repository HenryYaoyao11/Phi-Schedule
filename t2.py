from pptx import Presentation
from pptx.util import Inches

# 打开PPT文件
prs = Presentation("PPT.pptx")

# 获取第一张幻灯片
slide = prs.slides[0]

# 设置文本框的位置和尺寸
left, top, width, height = Inches(5549900), Inches(1819982), Inches(5425208), Inches(2853904)

# 遍历幻灯片中的所有形状
for shape in slide.shapes:
    # 判断形状是否为文本框
    if shape.has_text_frame:
        # 获取文本框的位置和尺寸
        text_box_left, text_box_top, text_box_width, text_box_height = shape.left, shape.top, shape.width, shape.height

        # 判断文本框的位置和尺寸是否与目标一致
        if (left, top, width, height) == (text_box_left, text_box_top, text_box_width, text_box_height):
            # 修改文本框内容为图片
            shape.text = ""
            picture = slide.shapes.add_picture("996.jpg", left, top, width, height)
            break

# 保存修改后的PPT文件
prs.save("修改后的PPT.pptx")
