from pptx import Presentation


# 打开PPT文件
ppt = Presentation("PPT.pptx")

# 获取第一张幻灯片
slide = ppt.slides[0]

# 设置文本框的位置和尺寸
left, top, width, height = 2809874, 3106057, 3286125, 1377329

# 遍历幻灯片中的所有形状
for shape in slide.shapes:
    # 判断形状是否为文本框
    if shape.has_text_frame:
        # 获取文本框的位置和尺寸
        text_box_left, text_box_top, text_box_width, text_box_height = shape.left, shape.top, shape.width, shape.height

        # 判断文本框的位置和尺寸是否与目标一致
        if (left, top, width, height) == (text_box_left, text_box_top, text_box_width, text_box_height):
            # 修改文本框内容
            shape.text = "xixi"
            break

# 设置文本框的位置和尺寸
left, top, width, height = 3288506, 2359820, 3064210, 746238

# 遍历幻灯片中的所有形状
for shape in slide.shapes:
    # 判断形状是否为文本框
    if shape.has_text_frame:
        # 获取文本框的位置和尺寸
        text_box_left, text_box_top, text_box_width, text_box_height = shape.left, shape.top, shape.width, shape.height

        # 判断文本框的位置和尺寸是否与目标一致
        if (left, top, width, height) == (text_box_left, text_box_top, text_box_width, text_box_height):
            # 修改文本框内容
            shape.text = "上课"
            break

# 设置文本框的位置和尺寸
left, top, width, height =11867892, 185766, 1613877,1054386

# 遍历幻灯片中的所有形状
for shape in slide.shapes:
    # 判断形状是否为文本框
    if shape.has_text_frame:
        # 获取文本框的位置和尺寸
        text_box_left, text_box_top, text_box_width, text_box_height = shape.left, shape.top, shape.width, shape.height

        # 判断文本框的位置和尺寸是否与目标一致
        if (left, top, width, height) == (text_box_left, text_box_top, text_box_width, text_box_height):
            # 修改文本框内容
            shape.text = "AA"
            break

# 设置文本框的位置和尺寸
left, top, width, height = 11492576, 1261632, 1613877, 1054386

# 遍历幻灯片中的所有形状
for shape in slide.shapes:
    # 判断形状是否为文本框
    if shape.has_text_frame:
        # 获取文本框的位置和尺寸
        text_box_left, text_box_top, text_box_width, text_box_height = shape.left, shape.top, shape.width, shape.height

        # 判断文本框的位置和尺寸是否与目标一致
        if (left, top, width, height) == (text_box_left, text_box_top, text_box_width, text_box_height):
            # 修改文本框内容
            shape.text = "222"
            break

# 设置文本框的位置和尺寸
left, top, width, height = 11118590, 2340967, 1613877, 1054386

# 遍历幻灯片中的所有形状
for shape in slide.shapes:
    # 判断形状是否为文本框
    if shape.has_text_frame:
        # 获取文本框的位置和尺寸
        text_box_left, text_box_top, text_box_width, text_box_height = shape.left, shape.top, shape.width, shape.height

        # 判断文本框的位置和尺寸是否与目标一致
        if (left, top, width, height) == (text_box_left, text_box_top, text_box_width, text_box_height):
            # 修改文本框内容
            shape.text = "333"
            break

# 设置文本框的位置和尺寸
left, top, width, height =10732921,3421380, 1613877, 1054386

# 遍历幻灯片中的所有形状
for shape in slide.shapes:
    # 判断形状是否为文本框
    if shape.has_text_frame:
        # 获取文本框的位置和尺寸
        text_box_left, text_box_top, text_box_width, text_box_height = shape.left, shape.top, shape.width, shape.height

        # 判断文本框的位置和尺寸是否与目标一致
        if (left, top, width, height) == (text_box_left, text_box_top, text_box_width, text_box_height):
            # 修改文本框内容
            shape.text = "444"
            break
# 设置文本框的位置和尺寸
left, top, width, height =10349752,4494173, 1613877, 1054386

# 遍历幻灯片中的所有形状
for shape in slide.shapes:
    # 判断形状是否为文本框
    if shape.has_text_frame:
        # 获取文本框的位置和尺寸
        text_box_left, text_box_top, text_box_width, text_box_height = shape.left, shape.top, shape.width, shape.height

        # 判断文本框的位置和尺寸是否与目标一致
        if (left, top, width, height) == (text_box_left, text_box_top, text_box_width, text_box_height):
            # 修改文本框内容
            shape.text = "555"
            break

# 设置文本框的位置和尺寸
left, top, width, height =9977901,5581128, 1613877, 1054386

# 遍历幻灯片中的所有形状
for shape in slide.shapes:
    # 判断形状是否为文本框
    if shape.has_text_frame:
        # 获取文本框的位置和尺寸
        text_box_left, text_box_top, text_box_width, text_box_height = shape.left, shape.top, shape.width, shape.height

        # 判断文本框的位置和尺寸是否与目标一致
        if (left, top, width, height) == (text_box_left, text_box_top, text_box_width, text_box_height):
            # 修改文本框内容
            shape.text = "666"
            break
# 设置文本框的位置和尺寸
left, top, width, height =9602816,6659737, 1613877, 1054386

# 遍历幻灯片中的所有形状
for shape in slide.shapes:
    # 判断形状是否为文本框
    if shape.has_text_frame:
        # 获取文本框的位置和尺寸
        text_box_left, text_box_top, text_box_width, text_box_height = shape.left, shape.top, shape.width, shape.height

        # 判断文本框的位置和尺寸是否与目标一致
        if (left, top, width, height) == (text_box_left, text_box_top, text_box_width, text_box_height):
            # 修改文本框内容
            shape.text = "777"
            break

# 设置文本框的位置和尺寸
left, top, width, height =5549900,1819982, 5425208, 2853904
from pptx import Presentation
from pptx.dml.color import RGBColor

# 遍历幻灯片中的所有形状
for shape in slide.shapes:
    # 判断形状是否为文本框
    if shape.has_text_frame:
        # 获取文本框的位置和尺寸
        text_box_left, text_box_top, text_box_width, text_box_height = shape.left, shape.top, shape.width, shape.height

        # 判断文本框的位置和尺寸是否与目标一致
        if (left, top, width, height) == (text_box_left, text_box_top, text_box_width, text_box_height):
            # 修改文本框内容为图片
            shape.text = ""
            picture = slide.shapes.add_picture("table\oimg\996.jpg", left, top, width, height)
            picture.rotation = shape.rotation
            shape.fill.solid()
            picture.fill.fore_color.rgb = shape.fill.fore_color.rgb
            break

# 保存修改后的PPT文件
ppt.save("修改后的PPT.pptx")
