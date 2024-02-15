from pptx import Presentation  
from pptx.util import Inches  
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE  
  
def extract_text_box_info(slide):  
    # 获取幻灯片中的所有形状  
    shapes = slide.shapes  
  
    # 遍历形状，找到文本框并提取编号、内容和位置  
    for shape in shapes:  
        if shape.has_text_frame:  
            text_frame = shape.text_frame  
            text_box_id = text_frame.text  # 提取文本框编号  
            text_content = text_frame.text  # 提取文本框内容（这里假设文本框编号同时也是内容）  
            left = shape.left  # 提取文本框左边缘位置  
            top = shape.top  # 提取文本框上边缘位置  
            width = shape.width  # 提取文本框宽度  
            height = shape.height  # 提取文本框高度  
            print(f"文本框编号：{text_box_id}")  
            print(f"文本框内容：{text_content}")  
            print(f"文本框位置：左={left}英寸, 上={top}英寸")  
            print(f"文本框尺寸：宽度={width}英寸, 高度={height}英寸")  
            print("-----------------------")  
  
# 打开PPT文件  
prs = Presentation('PPT.pptx')  
  
# 获取第一张幻灯片  
slide = prs.slides[0]  
  
# 提取文本框信息  
extract_text_box_info(slide)
