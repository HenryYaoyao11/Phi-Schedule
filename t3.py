from pptx import Presentation  
from pptx.util import Inches  
  
# 加载PPTX文件  
prs = Presentation('PPT2.pptx')  
  
# 导出幻灯片为JPEG格式的图片  
for slide in prs.slides:  
    for shape in slide.shapes:  
        if shape.has_text_frame:  
            # 获取文本框中的文本  
            text = shape.text_frame.text  
            # 创建与文本框相同尺寸的画布  
            canvas = shape.top_left_anchor.element.add_picture(f'{text}.jpg', width=Inches(shape.width), height=Inches(shape.height))  
            # 将画布保存为JPEG格式的图片  
            canvas.save_picture(f'{text}.jpg', format='JPEG')
